"""
generate_pptx.py
Generates a branded Cushman & Wakefield proposal PowerPoint.
Fetches a real property photo from Pexels based on property type.
"""

import os
import io
import urllib.request
import urllib.parse
import json
import tempfile
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Brand Colors ─────────────────────────────────────────────────────────────
CW_NAVY  = RGBColor(0x1A, 0x1A, 0x2E)
CW_RED   = RGBColor(0xE3, 0x27, 0x26)
CW_GRAY  = RGBColor(0x66, 0x66, 0x66)
CW_LIGHT = RGBColor(0xF4, 0xF4, 0xF4)
CW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CW_DARK  = RGBColor(0x2D, 0x2D, 0x2D)

def rgb(r, g, b): return RGBColor(r, g, b)


# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=12, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def add_footer(slide):
    add_rect(slide, 0, 5.35, 10, 0.275, CW_LIGHT)
    add_text(slide, "CONFIDENTIAL  |  © 2026 Cushman & Wakefield  |  All Rights Reserved",
             0.4, 5.35, 9.2, 0.275, font_size=8, color=CW_GRAY, align=PP_ALIGN.CENTER)


def add_header(slide):
    add_rect(slide, 0, 0, 10, 0.75, CW_NAVY)
    add_text(slide, "CUSHMAN & WAKEFIELD  |  ASSET SERVICES",
             0.4, 0.05, 9, 0.65, font_size=10, color=CW_WHITE)


def add_section_title(slide, title):
    add_text(slide, f"• {title}", 0.4, 0.88, 8, 0.35,
             font_size=20, bold=True, color=CW_NAVY)
    add_rect(slide, 0.4, 1.32, 9.2, 0.025, CW_RED)


# ── Pexels Image Fetcher ──────────────────────────────────────────────────────
PEXELS_QUERIES = {
    "Office Building":      "modern office building exterior",
    "Industrial Warehouse": "industrial warehouse building",
    "Retail Strip Center":  "retail shopping center storefront",
    "Mixed-Use Property":   "mixed use urban building",
    "Medical Office":       "medical office building",
    "Other":                "commercial real estate building",
}

def fetch_pexels_image(property_type, api_key):
    import requests
   
    query = PEXELS_QUERIES.get(property_type, "commercial building")
    try:
        headers = {"Authorization": api_key}
        response = requests.get(
            f"https://api.pexels.com/v1/search",
            headers=headers,
            params={"query": query, "per_page": 1, "orientation": "landscape"},
            timeout=10
        )
        data = response.json()
        if data["photos"]:
            img_url = data["photos"][0]["src"]["large"]
            img_response = requests.get(img_url, timeout=15)
            return img_response.content
    except Exception as e:
        print(f"Pexels fetch failed: {e}")
    return None


# ── Main Generator ────────────────────────────────────────────────────────────
def generate_proposal_pptx(client_name, property_address, property_type,
                            services, proposal_text, pexels_api_key,
                            headshot_path=None, output_path=None):

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]
    today = date.today().strftime("%B %d, %Y")

    # Fetch property image
    property_img_bytes = fetch_pexels_image(property_type, pexels_api_key)
    property_img_stream = io.BytesIO(property_img_bytes) if property_img_bytes else None

    # ── SLIDE 1: COVER ────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, 10, 5.625, CW_NAVY)
    add_rect(s1, 0, 0, 0.18, 5.625, CW_RED)
    add_rect(s1, 7.2, 0, 2.8, 5.625, rgb(0x2D, 0x2D, 0x2D))
    add_text(s1, "CUSHMAN &", 0.4, 1.0, 6, 0.55, font_size=36, bold=True, color=CW_WHITE)
    add_text(s1, "WAKEFIELD", 0.4, 1.55, 6, 0.55, font_size=36, bold=True, color=CW_RED)
    add_text(s1, "Asset Services", 0.4, 2.25, 6, 0.35, font_size=14, color=rgb(0xAA,0xAA,0xAA))
    add_text(s1, "PROPERTY SERVICES PROPOSAL", 0.4, 3.1, 9, 0.7, font_size=28, bold=True, color=CW_WHITE)
    add_text(s1, f"Prepared for: {client_name}", 0.4, 3.85, 9, 0.4, font_size=14, italic=True, color=rgb(0xCC,0xCC,0xCC))
    add_text(s1, today, 0.4, 5.1, 4, 0.3, font_size=11, color=rgb(0x99,0x99,0x99))
    add_text(s1, "MID-ATLANTIC\nREGION", 7.3, 2.3, 2.6, 0.9, font_size=11, color=rgb(0xAA,0xAA,0xAA), align=PP_ALIGN.CENTER)

    # ── SLIDE 2: PROPERTY SHOWCASE ────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    add_header(s2)

    if property_img_stream:
        s2.shapes.add_picture(property_img_stream, Inches(0.4), Inches(0.9), Inches(5.5), Inches(3.8))
    else:
        add_rect(s2, 0.4, 0.9, 5.5, 3.8, rgb(0x2C, 0x3E, 0x50))
        add_text(s2, f"[{property_type}]", 0.4, 2.5, 5.5, 0.6, font_size=16, bold=True, color=CW_WHITE, align=PP_ALIGN.CENTER)

    add_rect(s2, 0.4, 4.35, 5.5, 0.35, CW_NAVY)
    add_text(s2, property_type.upper(), 0.5, 4.37, 5.3, 0.3, font_size=11, bold=True, color=CW_WHITE, align=PP_ALIGN.CENTER)

    add_text(s2, "• PROPERTY DETAILS", 6.15, 0.9, 3.6, 0.35, font_size=16, bold=True, color=CW_NAVY)
    add_rect(s2, 6.15, 1.32, 3.6, 0.025, CW_RED)

    details = [
        ("CLIENT", client_name),
        ("ADDRESS", property_address),
        ("PROPERTY TYPE", property_type),
        ("SERVICES", ", ".join(services)),
    ]
    for i, (label, value) in enumerate(details):
        y = 1.5 + i * 0.78
        add_rect(s2, 6.15, y, 3.6, 0.68, CW_LIGHT)
        add_rect(s2, 6.15, y, 3.6, 0.06, CW_RED)
        add_text(s2, label, 6.25, y+0.08, 3.4, 0.2, font_size=8, bold=True, color=CW_RED)
        add_text(s2, value, 6.25, y+0.3, 3.4, 0.3, font_size=10, color=CW_NAVY)

    add_footer(s2)

    # ── SLIDE 3: PROPOSAL CONTENT ─────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    add_header(s3)
    add_section_title(s3, "PROPOSAL OVERVIEW")
    import re as re2
    clean = re2.sub(r'#{1,6}\s*', '', proposal_text)
    clean = re2.sub(r'\*\*(.*?)\*\*', r'\1', clean)
    clean = re2.sub(r'\*(.*?)\*', r'\1', clean)
    clean = re2.sub(r'---+', '', clean)
    clean = re2.sub(r'\n{3,}', '\n\n', clean).strip()
    preview = clean[:600]
    add_text(s3, preview, 0.4, 1.5, 9.2, 3.7, font_size=10, color=CW_DARK)
    add_footer(s3)

    # ── SLIDE 4: SCOPE OF SERVICES ────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    add_header(s4)
    add_section_title(s4, "SCOPE OF SERVICES")

    service_descs = {
        "Property Management": "Day-to-day operations, tenant relations, lease administration, vendor management, and 24/7 emergency response.",
        "Engineering Services": "Preventive maintenance, HVAC/electrical/plumbing systems, energy management, and capital project oversight.",
        "Janitorial Services": "Daily cleaning programs, green cleaning products, floor care, window washing, and emergency cleanup.",
        "Security Services": "24/7 security personnel, access control, CCTV monitoring, visitor management, and risk assessments.",
        "Accounting Services": "Full financial reporting, AP/AR, budget preparation, CAM reconciliations, and audit support.",
    }
    positions = [(0.4,1.5),(5.1,1.5),(0.4,3.0),(5.1,3.0),(0.4,4.25)]
    for i, svc in enumerate(services[:5]):
        x, y = positions[i]
        add_rect(s4, x, y, 4.5, 1.2, CW_WHITE)
        add_rect(s4, x, y, 0.08, 1.2, CW_RED)
        add_text(s4, svc.upper(), x+0.2, y+0.1, 4.2, 0.28, font_size=11, bold=True, color=CW_NAVY)
        desc = service_descs.get(svc, "Professional services tailored to your property needs.")
        add_text(s4, desc, x+0.2, y+0.42, 4.2, 0.65, font_size=9, color=CW_GRAY)

    add_footer(s4)

    # ── SLIDE 5: WHY C&W ──────────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    add_rect(s5, 0, 0, 10, 5.625, CW_NAVY)
    add_rect(s5, 0, 0, 0.18, 5.625, CW_RED)
    add_text(s5, "WHY CUSHMAN &", 0.4, 0.6, 9, 0.5, font_size=30, bold=True, color=CW_WHITE)
    add_text(s5, "WAKEFIELD", 0.4, 1.1, 9, 0.5, font_size=30, bold=True, color=CW_RED)

    reasons = [
        ("LOCAL EXPERTISE", "Deep knowledge of the Mid-Atlantic market across DC, Maryland, and Virginia."),
        ("PROVEN TRACK RECORD", "Decades managing commercial properties for institutional and private clients."),
        ("TECHNOLOGY-DRIVEN", "Industry-leading platforms for work orders, reporting, and tenant communication."),
        ("DEDICATED TEAM", "Single point of contact with full team support and 24/7 emergency response."),
    ]
    rpos = [(0.4,1.85),(5.1,1.85),(0.4,3.4),(5.1,3.4)]
    for i, (title, body) in enumerate(reasons):
        x, y = rpos[i]
        add_rect(s5, x, y, 4.4, 1.3, CW_WHITE)
        add_rect(s5, x, y, 4.4, 0.06, CW_RED)
        add_text(s5, title, x+0.15, y+0.1, 4.1, 0.28, font_size=10, bold=True, color=CW_RED)
        add_text(s5, body, x+0.15, y+0.42, 4.1, 0.7, font_size=9, color=CW_DARK)

    # ── SLIDE 6: NEXT STEPS ───────────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank)
    add_header(s6)
    add_section_title(s6, "NEXT STEPS")

    steps = [
        ("01", "REVIEW PROPOSAL", "Please review and contact us with any questions."),
        ("02", "INTRODUCTORY MEETING", "We welcome the opportunity to present our team in detail."),
        ("03", "FINALIZE SCOPE", "We will tailor our approach to your specific requirements."),
        ("04", "EXECUTE AGREEMENT", "We will prepare a management agreement for execution."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = 1.6 + i * 0.82
        add_rect(s6, 0.4, y, 0.55, 0.55, CW_RED)
        add_text(s6, num, 0.4, y, 0.55, 0.55, font_size=16, bold=True, color=CW_WHITE, align=PP_ALIGN.CENTER)
        add_text(s6, title, 1.1, y+0.02, 8, 0.25, font_size=11, bold=True, color=CW_NAVY)
        add_text(s6, desc, 1.1, y+0.3, 8, 0.25, font_size=10, color=CW_GRAY)

    add_footer(s6)

    # ── SLIDE 7: BIO ──────────────────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank)
    add_rect(s7, 0, 0, 10, 5.625, CW_NAVY)
    add_rect(s7, 0, 0, 0.18, 5.625, CW_RED)
    add_text(s7, "YOUR POINT OF CONTACT", 0.4, 0.35, 9, 0.35, font_size=11, color=rgb(0xAA,0xAA,0xAA))
    add_rect(s7, 0.4, 0.78, 9.2, 0.025, CW_RED)

    if headshot_path and os.path.exists(headshot_path):
        s7.shapes.add_picture(headshot_path, Inches(0.5), Inches(1.0), Inches(2.8), Inches(2.8))

    add_text(s7, "Alisha Kalsi", 3.55, 1.05, 6, 0.6, font_size=30, bold=True, color=CW_WHITE)
    add_text(s7, "Senior Administrator, Asset Services", 3.55, 1.65, 6, 0.32, font_size=14, bold=True, color=CW_RED)
    add_text(s7, "Mid-Atlantic Region  |  Cushman & Wakefield", 3.55, 1.97, 6, 0.28, font_size=12, color=rgb(0xAA,0xAA,0xAA))
    add_rect(s7, 3.55, 2.35, 5.8, 0.02, rgb(0x44,0x44,0x44))

    bio = ("Alisha Kalsi serves as Senior Administrator for Asset Services in the Mid-Atlantic "
           "region at Cushman & Wakefield, supporting the Asset Services team and senior managing "
           "directors across a diverse commercial real estate portfolio. In her role, Alisha coordinates "
           "proposal development, client communications, vendor relationships, and administrative "
           "operations spanning the DC, Maryland, and Virginia markets. She brings a detail-oriented, "
           "client-first approach to every engagement, ensuring seamless execution across all aspects "
           "of property services delivery.")
    add_text(s7, bio, 3.55, 2.48, 6.1, 2.2, font_size=10, color=rgb(0xCC,0xCC,0xCC))

    add_rect(s7, 0.4, 4.9, 9.2, 0.45, rgb(0x2D,0x2D,0x2D))
    add_text(s7, "Cushman & Wakefield  |  Mid-Atlantic Asset Services  |  Washington, DC",
             0.4, 4.9, 9.2, 0.45, font_size=10, color=rgb(0xCC,0xCC,0xCC), align=PP_ALIGN.CENTER)
    add_rect(s7, 0, 5.35, 10, 0.275, rgb(0x11,0x11,0x11))
    add_text(s7, "CONFIDENTIAL  |  © 2026 Cushman & Wakefield  |  All Rights Reserved",
             0.4, 5.35, 9.2, 0.275, font_size=8, color=CW_GRAY, align=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────────────
    if not output_path:
        output_path = tempfile.mktemp(suffix=".pptx")
    prs.save(output_path)
    return output_path
